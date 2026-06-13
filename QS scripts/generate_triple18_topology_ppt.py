# -*- coding: utf-8 -*-
"""
generate_triple18_topology_ppt.py
=================================
根据 triple_qs_topology18.m 的 18 种 effMat（仅 A1 行对 B1–B3 有 ±1/0），
生成与「三菌体系」示意相近的网络图，并写入一页一图的 PPTX。

依赖（在命令行执行）:
    pip install matplotlib python-pptx

用法:
    python generate_triple18_topology_ppt.py
    python generate_triple18_topology_ppt.py -o my_topologies.pptx

说明:
  - 与仿真枚举一致：跨群体 Hill 边只画 **A1→B1..B3**；A2、A3 仅作节点占位，无 A2/A3→B 虚线。
  - 各模块内固定示意：C→A（蓝实线）、C→B（橙实线）、B⊣C（橙点线负反馈）。
"""
from __future__ import annotations

import argparse
import io
import os
import sys

import matplotlib

matplotlib.use("Agg")
import matplotlib.pyplot as plt
import matplotlib.patches as mpatches
from matplotlib.patches import FancyArrowPatch
import numpy as np

try:
    from pptx import Presentation
    from pptx.util import Inches
except ImportError as e:
    print("请先安装: pip install python-pptx matplotlib", file=sys.stderr)
    raise SystemExit(1) from e


def topology18_data():
    """与 triple_qs_topology18.m 同步的 18 组 (3x3 matrix, 中文标签)。"""
    eff_mats = []
    labels = []
    k = 0
    for j in range(1, 4):
        k += 1
        M = np.zeros((3, 3), dtype=float)
        M[0, j - 1] = 1
        eff_mats.append(M)
        labels.append(f"单群体促进：A1→B{j} 促进")
    for j in range(1, 4):
        k += 1
        M = np.zeros((3, 3), dtype=float)
        M[0, j - 1] = -1
        eff_mats.append(M)
        labels.append(f"单群体抑制：A1→B{j} 抑制")
    pairs = [(0, 1), (0, 2), (1, 2)]  # (B1,B2) 等 0-based 列
    for a, b in pairs:
        k += 1
        M = np.zeros((3, 3), dtype=float)
        M[0, a] = 1
        M[0, b] = 1
        eff_mats.append(M)
        labels.append(f"双群体促进：A1→B{a+1},B{b+1}")
    for a, b in pairs:
        k += 1
        M = np.zeros((3, 3), dtype=float)
        M[0, a] = -1
        M[0, b] = -1
        eff_mats.append(M)
        labels.append(f"双群体抑制：A1→B{a+1},B{b+1}")
    eff_mats.append(np.array([[1, -1, 0], [0, 0, 0], [0, 0, 0]], dtype=float))
    labels.append("双群体异号：A1→B1促进 B2抑制")
    eff_mats.append(np.array([[1, 0, -1], [0, 0, 0], [0, 0, 0]], dtype=float))
    labels.append("双群体异号：A1→B1促进 B3抑制")
    eff_mats.append(np.array([[0, 1, -1], [0, 0, 0], [0, 0, 0]], dtype=float))
    labels.append("双群体异号：A1→B2促进 B3抑制")
    eff_mats.append(np.array([[1, 1, -1], [0, 0, 0], [0, 0, 0]], dtype=float))
    labels.append("三群体：B1B2促进 B3抑制")
    eff_mats.append(np.array([[1, -1, 1], [0, 0, 0], [0, 0, 0]], dtype=float))
    labels.append("三群体：B1B3促进 B2抑制")
    eff_mats.append(np.array([[-1, 1, 1], [0, 0, 0], [0, 0, 0]], dtype=float))
    labels.append("三群体：B2B3促进 B1抑制")
    assert len(eff_mats) == 18 and len(labels) == 18
    return eff_mats, labels


def _setup_chinese_font():
    matplotlib.rcParams["font.sans-serif"] = [
        "Microsoft YaHei",
        "SimHei",
        "PingFang SC",
        "Noto Sans CJK SC",
        "Arial Unicode MS",
        "DejaVu Sans",
    ]
    matplotlib.rcParams["axes.unicode_minus"] = False


def _node_positions():
    """三组菌大致三角分布；每组内 C / A / B 相对位置。"""
    ang = np.radians([90, 210, 330])
    R = 1.35
    centers = np.column_stack([R * np.cos(ang), R * np.sin(ang)])
    pos = {}
    labels_cab = [("C", "green"), ("A", "#2a6ebd"), ("B", "#e8943c")]
    # 每组：C 略靠重心，A 朝外，B 侧向
    offs = [
        np.array([[-0.22, -0.18], [0.0, 0.32], [0.22, -0.12]]),  # module 1 top
        np.array([[-0.2, 0.12], [0.28, -0.05], [0.0, -0.28]]),  # module 2
        np.array([[0.18, 0.1], [-0.28, -0.05], [0.0, -0.28]]),  # module 3
    ]
    for mi in range(3):
        c = centers[mi]
        for k, (letter, _) in enumerate(labels_cab):
            name = f"{letter}{mi + 1}"
            pos[name] = c + offs[mi][k]
    return pos


def _draw_arrow(ax, p0, p1, *, color, lw, ls="solid", rad=0.0, mutation_scale=14):
    p0 = np.asarray(p0, float)
    p1 = np.asarray(p1, float)
    style = f"arc3,rad={rad}"
    arr = FancyArrowPatch(
        p0,
        p1,
        arrowstyle="-|>",
        color=color,
        linewidth=lw,
        linestyle=ls,
        mutation_scale=mutation_scale,
        shrinkA=8,
        shrinkB=8,
        connectionstyle=style,
        zorder=1,
    )
    ax.add_patch(arr)


def _draw_inhibition(ax, p0, p1, *, color, lw, ls="dashed", rad=0.0):
    """虚线 + 终点 T 形平头（抑制）。"""
    p0 = np.asarray(p0, float)
    p1 = np.asarray(p1, float)
    v = p1 - p0
    nv = np.linalg.norm(v)
    if nv < 1e-9:
        return
    u = v / nv
    # 略缩进终点，留出 T 形位置
    tip = p1 - u * 0.08
    perp = np.array([-u[1], u[0]])
    bar_w = 0.06
    ax.plot([p0[0], tip[0]], [p0[1], tip[1]], color=color, lw=lw, ls=ls, zorder=1)
    ax.plot(
        [tip[0] - perp[0] * bar_w, tip[0] + perp[0] * bar_w],
        [tip[1] - perp[1] * bar_w, tip[1] + perp[1] * bar_w],
        color=color,
        lw=lw * 1.1,
        solid_capstyle="butt",
        zorder=2,
    )


def render_topology_figure(eff_row_a1: np.ndarray, title_cn: str, slide_idx: int) -> bytes:
    """返回 PNG 字节。"""
    pos = _node_positions()
    fig, ax = plt.subplots(figsize=(11, 8.2), dpi=140)
    ax.set_aspect("equal")
    ax.axis("off")
    ax.set_xlim(-2.1, 2.1)
    ax.set_ylim(-2.0, 2.0)

    # 节点
    node_colors = {"A": "#2a6ebd", "B": "#e8943c", "C": "#3d9e4a"}
    for name, xy in pos.items():
        letter = name[0]
        circ = plt.Circle(xy, 0.13, color=node_colors[letter], ec="white", lw=1.5, zorder=5)
        ax.add_patch(circ)
        ax.text(xy[0], xy[1], name, ha="center", va="center", color="white", fontsize=11, fontweight="bold", zorder=6)

    # 固定骨架：Ci→Ai, Ci→Bi；Bi⊣Ci
    for i in range(1, 4):
        Ci, Ai, Bi = pos[f"C{i}"], pos[f"A{i}"], pos[f"B{i}"]
        _draw_arrow(ax, Ci, Ai, color=node_colors["A"], lw=2.0, ls="solid", rad=0.05)
        _draw_arrow(ax, Ci, Bi, color=node_colors["B"], lw=2.0, ls="solid", rad=-0.05)
        _draw_inhibition(ax, Bi, Ci, color=node_colors["B"], lw=1.8, ls=(0, (2, 3)), rad=0.0)

    # 仅 A1→B*（与 effMat 第一行一致）
    p_a1 = pos["A1"]
    for j in range(3):
        w = eff_row_a1[j]
        if w == 0:
            continue
        p_b = pos[f"B{j + 1}"]
        rad = 0.12 * (j - 1)  # 轻微错开弧
        if w > 0:
            _draw_arrow(ax, p_a1, p_b, color=node_colors["A"], lw=2.0, ls="dashed", rad=rad, mutation_scale=16)
        else:
            _draw_inhibition(ax, p_a1, p_b, color=node_colors["A"], lw=2.0, ls="dashed", rad=rad)

    ax.text(0, 1.85, "三菌体系", ha="center", va="center", fontsize=20, fontweight="bold")
    ax.text(0, 1.58, f"H{slide_idx}  {title_cn}", ha="center", va="center", fontsize=11, color="#333333")

    leg = [
        mpatches.Patch(color=node_colors["C"], label="C：输入/群体"),
        mpatches.Patch(color=node_colors["A"], label="A：信号"),
        mpatches.Patch(color=node_colors["B"], label="B：产物"),
    ]
    ax.legend(handles=leg, loc="lower center", ncol=3, frameon=True, fontsize=9, bbox_to_anchor=(0.5, -0.02))

    buf = io.BytesIO()
    fig.savefig(buf, format="png", bbox_inches="tight", facecolor="white")
    plt.close(fig)
    buf.seek(0)
    return buf.read()


def build_pptx(out_path: str) -> None:
    eff_mats, labels = topology18_data()
    _setup_chinese_font()
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]

    title_slide = prs.slides.add_slide(prs.slide_layouts[0])
    title_slide.shapes.title.text = "三菌体系 · 18 种调控拓扑（与 triple_qs_topology18 一致）"
    sub = (
        "每页 H1–H18：仅 A1→B1–B3 跨边不同（蓝虚线：促进 / 抑制）。\n"
        "A2、A3 在模型中无对 B 的 Hill 跨边，故图中不画 A2/A3→B。\n"
        "各模块内：C→A、C→B 与 B⊣C 为示意性公共骨架。"
    )
    for ph in title_slide.placeholders:
        if ph.placeholder_format.idx == 1:
            ph.text = sub
            break

    for i in range(18):
        M = eff_mats[i]
        row = M[0, :]
        png_bytes = render_topology_figure(row, labels[i], i + 1)
        slide = prs.slides.add_slide(blank)
        pic_left = Inches(0.35)
        pic_top = Inches(0.55)
        pic_w = Inches(12.6)
        pic_h = Inches(6.5)
        stream = io.BytesIO(png_bytes)
        slide.shapes.add_picture(stream, pic_left, pic_top, width=pic_w)

    out_dir = os.path.dirname(os.path.abspath(out_path))
    if out_dir:
        os.makedirs(out_dir, exist_ok=True)
    prs.save(out_path)
    print(f"已写入: {out_path}")


def main():
    here = os.path.dirname(os.path.abspath(__file__))
    default_out = os.path.join(here, "triple18_topologies.pptx")
    ap = argparse.ArgumentParser(description="生成 18 种三菌拓扑示意 PPTX")
    ap.add_argument("-o", "--output", default=default_out, help="输出 .pptx 路径")
    args = ap.parse_args()
    build_pptx(os.path.abspath(args.output))


if __name__ == "__main__":
    main()
